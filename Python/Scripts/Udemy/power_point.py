# Como criar um arquivo power_point em Python
# ! Sempre importar collections e collections.abc
import collections
import collections.abc # noqa
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Criando uma nova apresentação
apresentacao = Presentation()

# Adicionando slides a apresentação

# ! Gerando todos os layouts possíveis
# Ordem dos layouts (0 = primeiro, 1 = segundo, ..., 9 = último)
# for c in range(9):
#     apresentacao.slides.add_slide(apresentacao.slide_layouts[c])

# ! Criando slide padrão (com título e subtítulo) = layouts[0]
# slide1 = apresentacao.slides.add_slide(apresentacao.slide_layouts[0])
# Definindo os espaços de digitação (título e subtítulo)
# # titulo = slide1.placeholders[0]
# titulo = slide1.shapes.title
# subtitulo = slide1.placeholders[1]
# # Escrevendo o texto
# titulo.text = 'oi'
# subtitulo.text = 'olá'

# ! Criando slide em branco = layout[6]
slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])

# ! Alterar a distância e o tamanho dos elementos
# Inches = distância em polegadas
# x = altura em relação ao slide
x = Inches(1)
# y = largura em relação ao slide
y = Inches(1)
# altura = altura da caixa de texto
altura = Inches(2)
# largura = largura da caixa de texto
largura = Inches(2)

# ! Criando uma caixa de texto
# Shapes = formas/formatos
caixa_texto = slide.shapes.add_textbox(x, y, largura, altura)

# Criando textos simples
caixa_texto.text = 'teste'

# ! Criando parágrafos
text_frame = caixa_texto.text_frame
# Definindo o parágrafo
paragrafo = text_frame.add_paragraph()
# Adicionando o texto
paragrafo.text = 'R$ 500'
# Negritando o texto
paragrafo.font.bold = True
# Mudando tamanho da fonte
paragrafo.font.size = Pt(30)

# ! Criando gráficos
# Categorias (eixo x)
produtos = ['pão', 'sal']
# Valores (eixo y)
vendas = [1, 2]

# Definindo o tamanho do gráfico
x = Inches(1)
y = Inches(1)
altura = Inches(5)
largura = Inches(3)

# Organizando o gráfico
dados_grafico = CategoryChartData()
# Definindo as categorias
dados_grafico.categories = produtos
# Definindo os valores
dados_grafico.add_series('Vendas', vendas)
# Escolhendo o tipo de gráfico = https://python-pptx.readthedocs.io/en/latest/dev/analysis/cht-chart-overview.html # noqa
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, largura, altura, dados_grafico) # noqa

# Salvando o arquivo (não esquecer a extensão do arquivo)
# ! OBS: Sempre fechar o arquivo antes de fazer modificações
apresentacao.save('meuteste2.pptx')

print('Powerpoint criado com sucesso.')
